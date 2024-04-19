from pptx import Presentation
from pptx.util import Inches
import os

def add_custom_slide(prs, layout, image_path, image_name):
    # Utilizza il layout fornito per aggiungere una nuova slide
    slide = prs.slides.add_slide(layout)

    # Aggiungi un titolo alla slide
    title = slide.shapes.title
    title.text = image_name

    # Aggiungi l'immagine
    left = Inches(1)
    top = Inches(1.5)
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(5.5))

    # Aggiungi un riquadro di testo con il nome dell'immagine
    left = Inches(7)
    top = Inches(1)
    width = Inches(2)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = image_name

def create_presentation_from_template(images_folder, template_path, output_path):
    # Carica il modello di presentazione
    prs = Presentation(template_path)

    #while len(prs.slides) > 1:
    #    xml_slides = prs.slides._sldIdLst  # Accesso al livello XML
    #    xml_slides.remove(xml_slides[-1])  # Rimuove l'ultimo slide dal fondo

    # Utilizza il layout della prima slide del modello
    first_slide_layout = prs.slide_layouts[0]

    # Trova tutte le immagini nella cartella specificata
    images = [img for img in os.listdir(images_folder) if img.endswith(('.png', '.jpg', '.jpeg'))]

    # Aggiungi il numero specificato di slide
    for i in range(len(images)):  # Assicurati di non superare il numero di immagini disponibili
        image_path = os.path.join(images_folder, images[i])
        add_custom_slide(prs, first_slide_layout, image_path, images[i].split('.')[0])

    # Salva la nuova presentazione
    prs.save(output_path)

# Percorso della cartella delle immagini, del modello e del file di output
images_folder = 'Images/'
template_path = './template.pptx'
output_path = '/mnt/c/Users/Marco/Desktop/your_presentation.pptx'

create_presentation_from_template(images_folder, template_path, output_path)
